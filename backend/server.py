from fastapi import FastAPI, APIRouter, UploadFile, File
from fastapi.responses import StreamingResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
import base64
import io
import uuid
import asyncio
from pathlib import Path
from pydantic import BaseModel, Field, ConfigDict
from typing import List, Optional
from datetime import datetime, timezone
import httpx
ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env', override=True)


def get_cors_origins() -> list[str]:
    raw_origins = os.environ.get("CORS_ORIGINS", "http://localhost:3000,http://127.0.0.1:3000")
    origins = [origin.strip() for origin in raw_origins.split(",") if origin.strip()]
    if "*" in origins:
        return ["http://localhost:3000", "http://127.0.0.1:3000"]
    return origins or ["http://localhost:3000", "http://127.0.0.1:3000"]

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

app = FastAPI()
api_router = APIRouter(prefix="/api")

try:
    mongo_url = os.environ.get('MONGO_URL', 'mongodb://localhost:27017')
    client = AsyncIOMotorClient(mongo_url)
    db = client[os.environ.get('DB_NAME', 'fairmount_loversai')]
    logger.info("MongoDB connected successfully")
except Exception as e:
    logger.warning(f"MongoDB connection failed: {e}. Moodboard features will be disabled.")
    client = None
    db = None

# Models
class GenerateRequest(BaseModel):
    prompt: str
    function_type: Optional[str] = None
    theme: Optional[str] = None
    space: Optional[str] = None
    venue_image: Optional[str] = None  # base64 - The actual venue (from angle/upload)
    design_image: Optional[str] = None  # base64 - Design inspiration (from template/upload)
    reference_image: Optional[str] = None  # base64 - Legacy support (treated as design_image)
    venue_image_url: Optional[str] = None
    design_image_url: Optional[str] = None
    high_quality: Optional[bool] = None
    variant_count: Optional[int] = None
    service_mode: Optional[bool] = None  # True → service overlay; skips decoration wrapper + "no people" policy

class MoodboardImage(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    session_id: str
    image_data: str  # base64
    prompt: str
    filters: dict = {}
    created_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())

class MoodboardSaveRequest(BaseModel):
    session_id: str
    image_data: str
    prompt: str
    filters: dict = {}

class DownloadRequest(BaseModel):
    images: List[dict]  # [{image_data: base64, prompt: str}]

SYSTEM_PROMPT = """\
### SYSTEM ROLE
You are "DecorVision AI" — an expert wedding design assistant specializing in spatial reasoning, decor integration, and photorealistic visualization. You help wedding designers merge client-chosen decor concepts with real venue photography while ensuring dimensional accuracy, aesthetic harmony, and practical feasibility.

### CORE OBJECTIVE
Analyze user-provided venue images + decor references → Generate dimensionally-aware, photorealistic visualizations showing how the decor integrates into the venue → Provide feasibility notes, placement guidance, and optimization suggestions.

### INPUT SPECIFICATION (Require from User)
┌─ Venue Data ───────────────────────────────┐
│ • Photos: 3+ angles (wide, mid, detail)    │
│ • Dimensions: L×W×H (ft/m) + ceiling height│
│ • Scale refs: door height, tile size, etc. │
│ • Fixed elements: pillars, stage, outlets  │
│ • Lighting: natural/artificial, color temp │
└────────────────────────────────────────────┘

┌─ Decor Design Data ────────────────────────┐
│ • Inspiration: mood boards, product links  │
│ • Items list: florals, drapery, lighting, │
│   furniture, centerpieces, signage         │
│ • Style keywords: "rustic", "glam", etc.   │
│ • Color palette: HEX/RGB + material prefs  │
│ • Budget tier: low/mid/high (affects density)│
└────────────────────────────────────────────┘

┌─ Client Context ───────────────────────────┐
│ • Guest count + seating layout needs       │
│ • Cultural requirements: mandap/chuppah   │
│ • Must-have vs. flexible elements          │
│ • Accessibility/safety constraints         │
└────────────────────────────────────────────┘

### PROCESSING WORKFLOW
1️⃣ SPATIAL ANALYSIS
   • Estimate scale using architectural cues or provided dims
   • Map focal zones: entrance, ceremony, reception, photo areas
   • Analyze lighting direction, intensity, color temperature
   • Identify traffic flow & sightline constraints

2️⃣ DECOR MAPPING & ADAPTATION
   • Match decor items to appropriate venue zones
   • Scale elements proportionally (e.g., arch height vs ceiling)
   • Adjust colors to complement venue palette + ambient light
   • Apply style-consistent textures/materials

3️⃣ FEASIBILITY CHECK & OPTIMIZATION
   • Flag conflicts: oversized items, crowded layouts, view blocks
   • Suggest alternatives if decor doesn't fit dims/budget
   • Recommend placement tweaks for visual balance + guest flow
   • Validate structural safety (weight limits, mounting points)

4️⃣ VISUALIZATION GENERATION
   • Create 3-5 photorealistic composite images:
     - Wide establishing shot
     - Detail close-ups (centerpiece, backdrop)
     - Guest-eye perspective (seated view)
   • Optional: 2D layout diagram with measurements
   • Maintain consistent shadows, perspective, lighting

### OUTPUT DELIVERABLES
✅ Visualizations:
   • 3-5 high-res rendered images (JPG/PNG)
   • Optional: 360° interactive view or short animation

✅ Documentation:
   • Annotated placement guide with dimensions
   • Itemized decor list with quantities + estimated costs
   • Feasibility report: what works, trade-offs, alternatives

✅ Communication:
   • Client-friendly explanation of design choices
   • Highlight "wow factor" elements + practical notes
   • Invite feedback for iterative refinement

### CONSTRAINTS & ETHICS
⚠️ Never invent dimensions — state assumptions clearly
⚠️ Avoid unsafe suggestions (heavy decor on weak structures)
⚠️ Respect cultural sensitivities + accessibility needs
⚠️ Disclose: AI visualizations are conceptual, not construction plans
⚠️ Prioritize sustainability: suggest reusable/rentable options

### STYLE & TONE GUIDELINES
•  Professional yet warm — you're collaborating with a designer
•  Use simple language for client-facing notes
•  Be specific: "Place 8ft floral arch 2ft left of entrance pillar"
•  Proactively suggest: "Consider uplighting here to enhance texture"

### ERROR HANDLING
If input is insufficient:
→ Request missing data politely with examples
→ Offer to proceed with estimated assumptions (flagged clearly)
→ Provide low-fidelity mockup first if high-res isn't feasible

### EXAMPLE USER PROMPT TEMPLATE
"I'm designing a wedding for [150 guests] at [outdoor garden venue].
Venue dims: [40ft x 60ft, 12ft ceiling clearance].
Style: ['enchanted forest' with gold accents].
Must-haves: [floral arch, hanging lanterns, round tables with greenery].
Budget: [mid-tier].
Attached: [venue_photos.zip], [decor_inspiration.pdf].
Please show me how this decor integrates into my space with sizing notes."
"""

SYSTEM_PROMPT_ADVANCED_V2 = """
You are an expert AI Prompt Engineer for photorealistic wedding decor visualization.

Goal:
- Analyze venue image + decor reference image + user request.
- Detect mode:
  - MODE A: element placement (single/few items)
  - MODE B: full venue decoration (entire space)
- Produce one final Flux-ready prompt that preserves venue structure while transferring decor vibe.

Mode rules:
- Specific single elements (mandap, arch, backdrop, aisle decor) -> MODE A
- Event-level requests (haldi, mehndi, sangeet, reception, full decoration) -> MODE B
- Ambiguous requests -> default MODE B

Prompt requirements:
- Keep walls, ceiling, floor materials, architecture, and perspective recognizable.
- Decoration is additive: place ON, hang FROM, stand AGAINST venue elements.
- Include concrete details: flower species, fabrics, lighting temperature, composition cues.
- No people, vehicles, or structural edits.
- Output format: line 1 contains MODE label, next a single paragraph prompt.
"""

def strip_env(value: Optional[str]) -> str:
    if not value:
        return ""
    return value.strip().strip('"').strip("'")


def detect_mode(user_prompt: str) -> str:
    """Mode A: element placement, Mode B: full venue decoration."""
    p = (user_prompt or "").lower().strip()

    mode_a_triggers = [
        "add mandap", "place mandap", "add arch", "place arch", "add backdrop",
        "place backdrop", "only florals", "just flowers", "only fabric",
        "just drapery", "aisle decor", "entrance decor", "stage setup", "add a backdrop",
    ]
    if any(trigger in p for trigger in mode_a_triggers):
        return "A"

    mode_b_triggers = [
        "decorate complete", "full decoration", "decorate the entire", "decorate entire",
        "decorate like this", "decorate hall", "haldi", "mehndi", "sangeet",
        "reception", "nikah", "wedding decoration", "transform", "fill the venue",
        "complete setup", "complete decoration", "same decoration", "this style",
        "this theme", "like this", "decorate for", "apply this", "make it look",
        "same vibe", "same feel", "full setup", "lavender", "purple theme",
        "yellow theme", "decorate with", "coral", "minimal decoration",
        "grand decoration", "heavy decoration", "opulent", "entire space",
        "whole hall", "whole venue",
    ]
    if any(trigger in p for trigger in mode_b_triggers):
        return "B"

    if any(word in p for word in ["decorate", "decoration", "setup", "theme"]):
        return "B"

    return "B"


def normalize_base64(image_value: Optional[str]) -> Optional[str]:
    if not image_value:
        return None
    value = image_value.strip()
    if value.startswith("data:") and "," in value:
        return value.split(",", 1)[1]
    return value


async def fetch_image_as_base64(url: str) -> Optional[str]:
    if not url:
        return None
    try:
        async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as fetch_client:
            response = await fetch_client.get(url)
            if response.status_code != 200:
                logger.warning(f"Image URL fetch failed ({response.status_code}): {url[:120]}")
                return None
            if not response.content:
                logger.warning(f"Image URL returned empty content: {url[:120]}")
                return None
            return base64.b64encode(response.content).decode("utf-8")
    except Exception as fetch_error:
        logger.warning(f"Image URL fetch exception: {fetch_error}")
        return None


def build_prompt(req: GenerateRequest) -> str:
    base_prompt = (req.prompt or "").strip()
    # User's specific request leads so the model prioritises it
    parts = [base_prompt if base_prompt else "Design the venue decor."]
    if req.function_type:
        parts.append(f"Event type: {req.function_type}")
    if req.theme:
        parts.append(f"Theme: {req.theme}")
    if req.space:
        parts.append(f"Venue space: {req.space} at Taj")
    parts.extend([
        "Output goal: one photorealistic transformed venue image suitable for client presentation.",
        "Hard constraints: preserve architecture, camera angle, and structural geometry.",
        "Decor policy: additive decor only; do not alter walls, floor, ceiling, pillars, doors, or windows.",
    ])
    if not req.service_mode:
        parts.append("Scene policy: no people, no text overlays, no logos, no watermarks.")
    return "\n".join(parts)


def build_service_overlay_prompt(user_prompt: str) -> str:
    """Prompt for the service overlay pass — adds staff to an already-decorated venue image."""
    return (
        "The VENUE IMAGE shows a fully decorated event venue. Your task is to add hospitality service staff to it.\n\n"
        "CRITICAL RULES:\n"
        "• Keep ALL existing decor, furniture, lighting, floral arrangements, ceiling installations, floor treatment, "
        "and architecture EXACTLY as they appear — do not change a single element of the background.\n"
        "• Only ADD people (waiters, butler staff, guests) composited naturally into the scene.\n"
        "• Match lighting, shadows, and colour cast of added staff to the existing venue lighting.\n"
        "• Maintain the same camera angle, perspective, and depth of field.\n\n"
        f"{user_prompt}"
    )


def build_gemini_prompt(user_text: str, context: str, mode: str, has_venue: bool, has_design: bool) -> str:
    """Build a concise, focused image-edit prompt for Gemini. user_text is the raw user request."""
    ctx_line = f"Event context: {context}\n" if context else ""

    if has_venue and has_design:
        if mode == "A":
            return (
                f"REQUEST: {user_text}\n{ctx_line}\n"
                "TASK: Edit the VENUE IMAGE by placing only the specific decor elements shown in the DECORATION REFERENCE IMAGE.\n"
                "RULES:\n"
                "• Keep venue architecture unchanged: walls, pillars, windows, doors, floor, ceiling, camera angle.\n"
                "• Match the reference's exact colors, floral species, fabrics, materials, and scale.\n"
                "• CEILING: add hanging elements from reference (chandeliers, floral canopy, draping, lanterns).\n"
                "• FLOOR: add floor elements from reference (rangoli, carpet, runner) with correct perspective.\n"
                "• Render photorealistically: correct shadows, depth, reflections.\n"
                "• No people. No text. No logos."
            )
        else:
            return (
                f"REQUEST: {user_text}\n{ctx_line}\n"
                "TASK: Fully decorate the VENUE IMAGE in the style of the DECORATION REFERENCE IMAGE.\n"
                "RULES:\n"
                "• CEILING: add ceiling decor from reference (floral canopy, chandeliers, pendants, fabric draping, fairy lights).\n"
                "• FLOOR: add floor treatment from reference (rangoli, floral carpet, aisle runner, decorative covering).\n"
                "• WALLS: add wall decor, backdrops, draping, floral installations, stage/altar setups from reference.\n"
                "• COLOR & STYLE: copy exact color palette, floral species, fabric textures, lighting warmth from reference.\n"
                "• Keep venue architecture intact: structural walls, ceiling, floor, pillars, windows, doors, camera angle.\n"
                "• The result must look as if the same decorator who styled the reference has decorated this exact venue.\n"
                "• Render photorealistically: correct shadows, reflections, depth of field, material textures.\n"
                "• No people. No text. No logos."
            )
    elif has_venue:
        return (
            f"REQUEST: {user_text}\n{ctx_line}\n"
            "TASK: Decorate the VENUE IMAGE with elegant luxury event decor.\n"
            "• CEILING: add hanging floral canopy, chandeliers, or fairy lights.\n"
            "• FLOOR: add petal rangoli, floral carpet, or decorative aisle runner.\n"
            "• WALLS: add floral installations or fabric draping.\n"
            "• Keep all architecture unchanged. Render photorealistically. No people. No text."
        )
    elif has_design:
        return (
            f"REQUEST: {user_text}\n{ctx_line}\n"
            "Generate a photorealistic decorated luxury event venue in the style of the DECORATION REFERENCE IMAGE. "
            "Include ceiling and floor decorations. No people. No text."
        )
    else:
        return (
            f"REQUEST: {user_text}\n{ctx_line}\n"
            "Generate a photorealistic beautifully decorated luxury event venue at Taj. "
            "Include elaborate ceiling floral canopy or chandeliers and decorative floor treatment. No people. No text."
        )


def guess_mime_from_base64(image_base64: str) -> str:
    snippet = (image_base64 or "")[:32]
    if snippet.startswith("iVBORw0KGgo"):
        return "image/png"
    if snippet.startswith("/9j/"):
        return "image/jpeg"
    if snippet.startswith("UklGR"):
        return "image/webp"
    return "image/jpeg"


def base64_to_data_uri(image_base64: str) -> str:
    mime = guess_mime_from_base64(image_base64)
    return f"data:{mime};base64,{image_base64}"


def extract_message_text(content) -> str:
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, dict) and item.get("type") == "text":
                parts.append(item.get("text", ""))
        return "\n".join(part for part in parts if part).strip()
    return ""


def extract_final_prompt(raw_output: str, mode: str) -> str:
    text = (raw_output or "").strip()
    if not text:
        return ""

    lines = text.split("\n")
    cleaned = []
    for line in lines:
        s = line.strip().upper()
        if s in {
            "MODE A", "MODE B", "MODE: A", "MODE: B", "**MODE A**", "**MODE B**",
            "DETECTED MODE: A", "DETECTED MODE: B", "A", "B"
        }:
            continue
        cleaned.append(line)
    text = "\n".join(cleaned).strip()

    markers = [
        "**FINAL PROMPT:**", "FINAL PROMPT:", "**FINAL PROMPT**", "FINAL PROMPT",
        "**PROMPT:**", "PROMPT:"
    ]
    for marker in markers:
        idx = text.find(marker)
        if idx != -1:
            candidate = text[idx + len(marker):].strip().strip("-").strip()
            if candidate.startswith('"') and candidate.endswith('"'):
                candidate = candidate[1:-1].strip()
            if len(candidate.split()) > 50:
                return candidate

    for starter in [
        "A photorealistic photograph of", "A photorealistic photo of", "A photorealistic image of"
    ]:
        idx = text.find(starter)
        if idx != -1:
            candidate = text[idx:].strip()
            if len(candidate.split()) > 80:
                return candidate

    longest = ""
    for para in text.split("\n\n"):
        p = para.strip()
        if p and len(p.split()) > len(longest.split()) and p[0] not in "#*-=":
            longest = p
    if len(longest.split()) > 80:
        return longest

    min_wc = 280 if mode == "A" else 340
    prose_lines = []
    for line in text.split("\n"):
        s = line.strip()
        if not s:
            continue
        if s[0] in "#*-=":
            continue
        if ":" in s and len(s.split()) < 8:
            continue
        if len(s.split()) > 8:
            prose_lines.append(s)
    combined = " ".join(prose_lines).strip()
    if len(combined.split()) > min_wc:
        return combined

    return text


def build_analysis_instruction(mode: str, user_prompt: str) -> str:
    if mode == "B":
        return f"""Analyze both images for FULL VENUE DECORATION.
IMAGE 1 (VENUE): room type, floor material/pattern, walls/panels/trim, ceiling structure, lighting, perspective, top venue identity features.
IMAGE 2 (DECOR): event style, color palette, ceiling decor, wall decor, floor treatment, focal point, seating, florals (species), density, distribution.
USER REQUEST: \"{user_prompt}\"
Return concise structured analysis."""
    return f"""Analyze both images for ELEMENT PLACEMENT.
IMAGE 1 (VENUE): floor, walls, ceiling, lighting, perspective, top venue identity features.
IMAGE 2 (DECOR): decor element types, approximate dimensions, floral details, fabrics, top colors.
USER REQUEST: \"{user_prompt}\"
Return concise structured analysis."""


async def generate_prompt_with_groq(
    user_prompt: str,
    venue_image_base64: str,
    decor_image_base64: str,
    groq_api_key: str,
    groq_model: str,
) -> tuple[str, str, str]:
    mode = detect_mode(user_prompt)
    venue_data_uri = base64_to_data_uri(venue_image_base64)
    decor_data_uri = base64_to_data_uri(decor_image_base64)

    headers = {
        "Authorization": f"Bearer {groq_api_key}",
        "Content-Type": "application/json",
    }
    groq_url = "https://api.groq.com/openai/v1/chat/completions"

    analysis_instruction = build_analysis_instruction(mode, user_prompt)

    async with httpx.AsyncClient(timeout=90.0) as groq_client:
        analysis_payload = {
            "model": groq_model,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": analysis_instruction},
                        {"type": "image_url", "image_url": {"url": venue_data_uri}},
                        {"type": "image_url", "image_url": {"url": decor_data_uri}},
                    ],
                }
            ],
            "temperature": 0.05,
            "max_tokens": 2000,
        }
        analysis_resp = await groq_client.post(groq_url, headers=headers, json=analysis_payload)
        analysis_resp.raise_for_status()
        analysis_json = analysis_resp.json()
        analysis_content = analysis_json.get("choices", [{}])[0].get("message", {}).get("content", "")
        analysis = extract_message_text(analysis_content)

        prompt_payload = {
            "model": groq_model,
            "messages": [
                {"role": "system", "content": SYSTEM_PROMPT_ADVANCED_V2},
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": f"""Analysis:
---
{analysis}
---
User: \"{user_prompt}\"
Output MODE on line 1, then the prompt paragraph. Nothing else.""",
                        },
                        {"type": "image_url", "image_url": {"url": venue_data_uri}},
                        {"type": "image_url", "image_url": {"url": decor_data_uri}},
                    ],
                },
            ],
            "temperature": 0.2,
            "max_tokens": 3000,
        }
        prompt_resp = await groq_client.post(groq_url, headers=headers, json=prompt_payload)
        prompt_resp.raise_for_status()
        prompt_json = prompt_resp.json()
        raw_content = prompt_json.get("choices", [{}])[0].get("message", {}).get("content", "")
        raw_output = extract_message_text(raw_content)

        final_prompt = extract_final_prompt(raw_output, mode)
        min_words = 280 if mode == "A" else 340

        if len(final_prompt.split()) < min_words:
            retry_payload = {
                "model": groq_model,
                "messages": [
                    {"role": "system", "content": SYSTEM_PROMPT_ADVANCED_V2},
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": (
                                    f"Analysis:\n{analysis}\n\n"
                                    f"User: \"{user_prompt}\"\n\n"
                                    f"MODE + prompt only. {'380-480' if mode == 'B' else '320-420'} words. "
                                    "Start with 'A photorealistic photograph of...'"
                                ),
                            },
                            {"type": "image_url", "image_url": {"url": venue_data_uri}},
                            {"type": "image_url", "image_url": {"url": decor_data_uri}},
                        ],
                    },
                ],
                "temperature": 0.2,
                "max_tokens": 2500,
            }
            try:
                retry_resp = await groq_client.post(groq_url, headers=headers, json=retry_payload)
                retry_resp.raise_for_status()
                retry_json = retry_resp.json()
                retry_content = retry_json.get("choices", [{}])[0].get("message", {}).get("content", "")
                retry_output = extract_message_text(retry_content)
                retry_prompt = extract_final_prompt(retry_output, mode)
                if len(retry_prompt.split()) > len(final_prompt.split()):
                    final_prompt = retry_prompt
            except Exception as retry_error:
                logger.warning(f"Groq retry prompt enhancement failed: {retry_error}")

    if not final_prompt or len(final_prompt.split()) < 50:
        final_prompt = user_prompt

    return final_prompt, analysis, mode


def resolve_result_image_url(data: dict) -> Optional[str]:
    if not isinstance(data, dict):
        return None
    for key in ("sample", "url", "image_url", "output"):
        val = data.get(key)
        if isinstance(val, str) and val.startswith("http"):
            return val
    result = data.get("result")
    if isinstance(result, dict):
        for key in ("sample", "url", "image_url", "output"):
            val = result.get(key)
            if isinstance(val, str) and val.startswith("http"):
                return val
    elif isinstance(result, str) and result.startswith("http"):
        return result
    return None


def guess_mime_from_binary(data: bytes) -> str:
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return "image/png"
    if data.startswith(b"\xff\xd8\xff"):
        return "image/jpeg"
    if data.startswith(b"RIFF") and b"WEBP" in data[:16]:
        return "image/webp"
    return "image/png"


def build_nano_banana_contents(prompt: str, venue_image: Optional[str], design_image: Optional[str]) -> list:
    parts = []

    if venue_image and design_image:
        parts.append({"text": "VENUE IMAGE — this is the real venue; preserve its architecture, camera angle, perspective, walls, floor, ceiling, pillars, doors, and windows exactly:"})
        parts.append({"inlineData": {"mimeType": guess_mime_from_base64(venue_image), "data": venue_image}})
        parts.append({"text": "DECORATION REFERENCE IMAGE — copy this decor style, color palette, flowers, fabrics, lighting, and atmosphere faithfully:"})
        parts.append({"inlineData": {"mimeType": guess_mime_from_base64(design_image), "data": design_image}})
    elif venue_image:
        parts.append({"text": "VENUE IMAGE — decorate this venue:"})
        parts.append({"inlineData": {"mimeType": guess_mime_from_base64(venue_image), "data": venue_image}})
    elif design_image:
        parts.append({"text": "DECORATION REFERENCE IMAGE:"})
        parts.append({"inlineData": {"mimeType": guess_mime_from_base64(design_image), "data": design_image}})

    parts.append({"text": prompt})
    return [{"role": "user", "parts": parts}]


def extract_nano_banana_image_and_text(data: dict) -> tuple[Optional[str], Optional[str], Optional[str]]:
    """Returns (image_base64, mime_type, text_message)."""
    if not isinstance(data, dict):
        return None, None, None

    candidates = data.get("candidates") or []
    text_fragments = []
    for candidate in candidates:
        content = (candidate or {}).get("content") or {}
        for part in content.get("parts") or []:
            if isinstance(part, dict):
                inline = part.get("inlineData") or part.get("inline_data")
                if isinstance(inline, dict) and inline.get("data"):
                    return inline.get("data"), inline.get("mimeType") or inline.get("mime_type") or "image/png", None
                if part.get("text"):
                    text_fragments.append(str(part.get("text")))

    return None, None, "\n".join(text_fragments).strip() or None


async def generate_with_nano_banana(
    prompt: str,
    api_key: str,
    model: str,
    venue_image: Optional[str],
    design_image: Optional[str],
) -> dict:
    payload = {
        "contents": build_nano_banana_contents(prompt, venue_image, design_image),
        "generationConfig": {
            "responseModalities": ["TEXT", "IMAGE"],
        },
    }

    # Try the configured model first, then fall back to known image generation models.
    # gemini-2.5-flash-image is confirmed working; others are kept as safety nets.
    candidate_models = [
        model,
        "gemini-2.5-flash-image",
        "gemini-2.0-flash-preview-image-generation",
        "gemini-2.0-flash-exp-image-generation",
    ]

    # Preserve order while removing duplicates.
    seen: set[str] = set()
    models_to_try: list[str] = []
    for candidate in candidate_models:
        c = (candidate or "").strip()
        if c and c not in seen:
            seen.add(c)
            models_to_try.append(c)

    async with httpx.AsyncClient(timeout=120.0) as nb_client:
        for idx, model_name in enumerate(models_to_try):
            url = (
                f"https://generativelanguage.googleapis.com/v1beta/models/"
                f"{model_name}:generateContent?key={api_key}"
            )
            retryable_attempts = 0

            while True:
                response = await nb_client.post(
                    url, json=payload, headers={"Content-Type": "application/json"}
                )
                body_text = response.text.strip()

                if response.status_code in {429, 500, 502, 503, 504} and retryable_attempts < 2:
                    retryable_attempts += 1
                    logger.warning(
                        f"Nano Banana transient error ({model_name}), retry {retryable_attempts}/2. "
                        f"Status={response.status_code}, body={body_text[:300]}"
                    )
                    await asyncio.sleep(2 * retryable_attempts)
                    continue

                break

            if response.status_code == 200:
                try:
                    data = response.json()
                except Exception:
                    return {
                        "success": False,
                        "error": "Nano Banana invalid JSON response",
                        "description": body_text[:1000],
                    }

                image_b64, mime_type, text_message = extract_nano_banana_image_and_text(data)
                if not image_b64:
                    # Model responded but returned no image — don't fall through to next model,
                    # return the error so the caller can surface it.
                    return {
                        "success": False,
                        "error": "Nano Banana returned no image",
                        "description": text_message or body_text[:1000] or "No image found in response",
                    }

                return {
                    "success": True,
                    "image_data": image_b64,
                    "mime_type": mime_type or "image/png",
                    "description": "Design generated successfully",
                    "provider": "nano_banana",
                    "model": model_name,
                }

            has_next = idx < len(models_to_try) - 1
            is_model_not_found = response.status_code == 404 and "not found" in body_text.lower()
            is_transient = response.status_code in {429, 500, 502, 503, 504}

            if has_next and (is_model_not_found or is_transient):
                logger.warning(
                    f"Nano Banana model attempt failed ({model_name}), trying next. "
                    f"Status={response.status_code}, body={body_text[:300]}"
                )
                continue

            logger.error(f"Nano Banana API error {response.status_code}: {body_text[:1000]}")
            return {
                "success": False,
                "error": f"Nano Banana API error: {response.status_code}",
                "description": body_text[:1000] or "Request failed",
            }

    return {
        "success": False,
        "error": "Nano Banana model resolution failed",
        "description": "No compatible Nano Banana model could be used.",
    }


async def generate_flux_variant_from_payload(
    target_url: str,
    api_payload: dict,
    flux_api_key: str,
    max_attempts: int = 100,
) -> dict:
    """Generate a single FLUX image variant using the BFL API with polling."""
    async with httpx.AsyncClient(timeout=90.0) as submit_client:
        response = await submit_client.post(
            target_url,
            json=api_payload,
            headers={
                "x-key": flux_api_key,
                "Content-Type": "application/json",
            },
        )

    if response.status_code not in [200, 201, 202]:
        body = response.text.strip()
        return {
            "success": False,
            "error": f"Variant submit failed: {response.status_code}",
            "description": body[:500] if body else "Submit failed",
        }

    try:
        data = response.json()
    except Exception:
        return {
            "success": False,
            "error": "Variant submit invalid response",
            "description": response.text[:500],
        }

    task_id = data.get("id", "")
    polling_url = data.get("polling_url", f"https://api.bfl.ai/v1/get_result?id={task_id}")
    result_url = resolve_result_image_url(data)

    if not result_url and not task_id:
        return {
            "success": False,
            "error": "Variant missing task ID or image URL",
            "description": "Flux response did not include polling info",
        }

    if not result_url:
        async with httpx.AsyncClient(timeout=15.0) as poll_client:
            for _ in range(max_attempts):
                await asyncio.sleep(3)
                poll_response = await poll_client.get(
                    polling_url,
                    headers={"x-key": flux_api_key, "Content-Type": "application/json"},
                )
                if poll_response.status_code != 200:
                    continue
                poll_data = poll_response.json()
                status = str(poll_data.get("status", "unknown")).lower()
                if status in ["ready", "succeeded", "complete", "completed"]:
                    result_url = resolve_result_image_url(poll_data)
                    if result_url:
                        break
                    return {
                        "success": False,
                        "error": "Variant result missing image",
                        "description": str(poll_data)[:200],
                    }
                if status in ["error", "failed", "content moderated", "request moderated"]:
                    return {
                        "success": False,
                        "error": "Variant generation failed",
                        "description": str(poll_data)[:200],
                    }

    if not result_url:
        return {
            "success": False,
            "error": "Variant timeout",
            "description": "Variant generation timed out",
        }

    async with httpx.AsyncClient(timeout=60.0) as dl_client:
        img_response = await dl_client.get(result_url)
    if img_response.status_code != 200:
        return {
            "success": False,
            "error": "Variant download failed",
            "description": f"Status: {img_response.status_code}",
        }

    img_base64 = base64.b64encode(img_response.content).decode("utf-8")
    return {
        "success": True,
        "image_data": img_base64,
        "mime_type": "image/png",
        "task_id": task_id,
    }


@api_router.get("/")
async def root():
    return {"message": "Taj x LoversAI API"}

@api_router.post("/generate")
async def generate_image(req: GenerateRequest):
    try:
        # Use raw user text — avoid double-wrapping through build_prompt
        raw_user_text = (req.prompt or "").strip() or "Design the venue decor."
        mode = detect_mode(raw_user_text)

        # Resolve images from base64 or URL first — needed to build the correct prompt
        venue_image = normalize_base64(req.venue_image)
        design_image = normalize_base64(req.design_image)

        if not venue_image and req.venue_image_url:
            logger.info("Venue base64 missing, fetching from venue_image_url...")
            venue_image = await fetch_image_as_base64(req.venue_image_url)

        if not design_image and req.design_image_url:
            logger.info("Design base64 missing, fetching from design_image_url...")
            design_image = await fetch_image_as_base64(req.design_image_url)

        if req.reference_image and not design_image:
            design_image = normalize_base64(req.reference_image)
            logger.info("Using legacy reference_image as design_image")

        has_venue = bool(venue_image)
        has_design = bool(design_image)

        # Short context line (event type + space) — appended inside build_gemini_prompt
        context_parts = []
        if req.function_type:
            context_parts.append(req.function_type)
        if req.space:
            context_parts.append(f"{req.space} at Taj")
        context_str = " | ".join(context_parts)

        # For service overlay, skip decoration wrapper — pass directly so waiters/staff are generated
        if req.service_mode:
            gemini_prompt = build_service_overlay_prompt(raw_user_text)
        else:
            gemini_prompt = build_gemini_prompt(raw_user_text, context_str, mode, has_venue, has_design)

        requested_variants = req.variant_count if req.variant_count is not None else (3 if req.high_quality else 1)
        variant_count = max(1, min(3, int(requested_variants or 1)))
        requested_input_mode = (
            "dual_image" if (has_venue and has_design)
            else "venue_only" if has_venue
            else "design_only" if has_design
            else "text_only"
        )

        logger.info(f"=== API RECEIVED ===")
        logger.info(f"prompt: {raw_user_text[:80]}...")
        logger.info(f"function_type: {req.function_type}")
        logger.info(f"space: {req.space}")
        logger.info(f"venue_image provided: {has_venue}, length: {len(venue_image) if venue_image else 0}")
        logger.info(f"design_image provided: {has_design}, length: {len(design_image) if design_image else 0}")
        logger.info(f"reference_image provided: {bool(req.reference_image)}")
        logger.info(f"venue_image_url provided: {bool(req.venue_image_url)}")
        logger.info(f"design_image_url provided: {bool(req.design_image_url)}")
        logger.info(f"variant_count: {variant_count}, mode: {mode}, input_mode: {requested_input_mode}")
        logger.info(f"gemini_prompt (first 120 chars): {gemini_prompt[:120]}")

        # =========================================================
        # Provider: Nano Banana (Gemini image generation)
        # =========================================================
        nano_banana_api_key = (
            strip_env(os.environ.get("NANO_BANANA_API_KEY", ""))
            or strip_env(os.environ.get("GEMINI_API_KEY", ""))
        )
        nano_banana_model = strip_env(os.environ.get("NANO_BANANA_MODEL", "gemini-2.5-flash-image"))

        if nano_banana_api_key:
            logger.info(f"Using Nano Banana (Gemini) for image generation — model: {nano_banana_model}")
            nano_tasks = [
                generate_with_nano_banana(
                    prompt=gemini_prompt,
                    api_key=nano_banana_api_key,
                    model=nano_banana_model,
                    venue_image=venue_image,
                    design_image=design_image,
                )
                for _ in range(variant_count)
            ]
            nano_results = await asyncio.gather(*nano_tasks)
            nano_variants = []
            nano_failure = None

            for nano_result in nano_results:
                if nano_result.get("success"):
                    nano_variants.append({
                        "image_data": nano_result.get("image_data"),
                        "mime_type": nano_result.get("mime_type", "image/png"),
                        "provider": "nano_banana",
                        "model": nano_result.get("model", nano_banana_model),
                    })
                else:
                    nano_failure = nano_failure or nano_result

            if nano_variants:
                primary = nano_variants[0]
                return {
                    "success": True,
                    "image_data": primary.get("image_data"),
                    "mime_type": primary.get("mime_type", "image/png"),
                    "description": "Design generated successfully",
                    "provider": "nano_banana",
                    "model": primary.get("model", nano_banana_model),
                    "input_mode": requested_input_mode,
                    "effective_input_mode": requested_input_mode,
                    "mode": "element" if mode == "A" else "full_decoration",
                    "prompt_source": "nano_banana",
                    "prompt_words": len(gemini_prompt.split()),
                    "variants": nano_variants,
                    "variants_requested": variant_count,
                    "variants_generated": len(nano_variants),
                }

            failure_desc = str(nano_failure.get("description", "")) if nano_failure else ""
            failure_error = str(nano_failure.get("error", "")) if nano_failure else ""
            logger.error(f"Nano Banana generation failed: {failure_error} — {failure_desc[:200]}")
            return {
                "success": False,
                "error": failure_error or "Nano Banana generation failed",
                "description": failure_desc or "No image was produced.",
            }

        return {
            "success": False,
            "error": "No image generation provider configured",
            "description": "Set NANO_BANANA_API_KEY or GEMINI_API_KEY in the .env file.",
        }

    except httpx.TimeoutException:
        logger.error("Image generation API timeout")
        return {"success": False, "error": "Request timed out", "description": "Try again later"}
    except Exception as e:
        logger.error(f"Generation error: {str(e)}")
        return {"success": False, "error": str(e)}

@api_router.post("/moodboard/save")
async def save_to_moodboard(req: MoodboardSaveRequest):
    if db is None:
        return {"success": False, "error": "Database not available", "id": None}
    doc = {
        "id": str(uuid.uuid4()),
        "session_id": req.session_id,
        "image_data": req.image_data,
        "prompt": req.prompt,
        "filters": req.filters,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    await db.moodboard_images.insert_one(doc)
    return {"success": True, "id": doc["id"]}

@api_router.get("/moodboard/{session_id}")
async def get_moodboard(session_id: str):
    if db is None:
        return {"images": [], "error": "Database not available"}
    images = await db.moodboard_images.find(
        {"session_id": session_id}, {"_id": 0}
    ).to_list(100)
    return {"images": images}

@api_router.post("/moodboard/download-pdf")
async def download_pdf(req: DownloadRequest):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas as pdf_canvas
    from reportlab.lib.utils import ImageReader

    buffer = io.BytesIO()
    c = pdf_canvas.Canvas(buffer, pagesize=A4)
    width, height = A4

    c.setFont("Helvetica-Bold", 28)
    c.drawCentredString(width / 2, height - 2 * inch, "Taj x LoversAI")
    c.setFont("Helvetica", 16)
    c.drawCentredString(width / 2, height - 2.6 * inch, "Moodboard Collection")
    c.setFont("Helvetica-Oblique", 11)
    c.drawCentredString(width / 2, height - 3.2 * inch, datetime.now(timezone.utc).strftime("%B %d, %Y"))
    c.showPage()

    for i, img_item in enumerate(req.images):
        try:
            img_bytes = base64.b64decode(img_item["image_data"])
            img_reader = ImageReader(io.BytesIO(img_bytes))

            img_w = width - 2 * inch
            img_h = img_w * 0.65
            x = inch
            y = height - img_h - 1.5 * inch

            c.drawImage(img_reader, x, y, width=img_w, height=img_h, preserveAspectRatio=True)

            c.setFont("Helvetica", 10)
            prompt_text = img_item.get("prompt", "")[:80]
            c.drawString(inch, y - 0.3 * inch, f"Design {i + 1}: {prompt_text}")
            c.showPage()
        except Exception as e:
            logger.error(f"PDF image error: {e}")
            continue

    c.save()
    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=moodboard.pdf"}
    )

@api_router.post("/moodboard/download-ppt")
async def download_ppt(req: DownloadRequest):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Taj x LoversAI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Moodboard Collection"
    p2.font.size = Pt(24)
    p2.alignment = PP_ALIGN.CENTER

    for i, img_item in enumerate(req.images):
        try:
            img_bytes = base64.b64decode(img_item["image_data"])
            img_stream = io.BytesIO(img_bytes)

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(img_stream, Inches(1), Inches(0.5), Inches(11.333), Inches(6))

            txBox = slide.shapes.add_textbox(Inches(1), Inches(6.7), Inches(11), Inches(0.6))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"Design {i + 1}: {img_item.get('prompt', '')[:60]}"
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.LEFT
        except Exception as e:
            logger.error(f"PPT image error: {e}")
            continue

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=moodboard.pptx"}
    )

import json

TEMPLATES_DIR = ROOT_DIR / "data" / "templates"
TEMPLATES_CONFIG = ROOT_DIR / "data" / "templates.json"

@api_router.get("/templates")
async def list_templates():
    try:
        if TEMPLATES_CONFIG.exists():
            with open(TEMPLATES_CONFIG, "r") as f:
                templates = json.load(f)
            for t in templates:
                t["available"] = (TEMPLATES_DIR / t["filename"]).exists()
            return {"templates": templates}
        return {"templates": []}
    except Exception as e:
        logger.error(f"Templates error: {e}")
        return {"templates": []}

@api_router.get("/templates/download/{filename}")
async def download_template(filename: str):
    filepath = TEMPLATES_DIR / filename
    if not filepath.exists():
        return {"error": "Template file not found"}
    return StreamingResponse(
        open(filepath, "rb"),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=get_cors_origins(),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("shutdown")
async def shutdown_db_client():
    if client is not None:
        client.close()